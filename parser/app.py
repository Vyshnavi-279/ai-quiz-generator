import os
import sys
import traceback
from io import BytesIO

from flask import Flask, jsonify, request
from flask_cors import CORS
from pptx import Presentation

app = Flask(__name__)
CORS(app)

MAX_FILE_SIZE = 25 * 1024 * 1024  # 25 MB


def extract_slide_text(slide, slide_index):
    """
    Extract text content from a single slide.
    Returns a dict with slideNumber, title, content, and wordCount.
    """
    all_text_parts = []
    first_text = None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                all_text_parts.append(text)
                if first_text is None:
                    first_text = text

    content = ' '.join(all_text_parts)
    title = first_text if first_text else f'Slide {slide_index + 1}'
    word_count = len(content.split())

    return {
        'slideNumber': slide_index + 1,
        'title': title,
        'content': content,
        'wordCount': word_count,
    }


@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'ok'})


@app.route('/parse', methods=['POST'])
def parse_pptx():
    # Check if a file was provided
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided. Please attach a .pptx file with field name "file".'}), 400

    file = request.files['file']

    if file.filename == '':
        return jsonify({'error': 'No file selected.'}), 400

    # Validate file extension
    if not file.filename.lower().endswith('.pptx'):
        return jsonify({'error': 'Invalid file type. Only .pptx files are allowed.'}), 400

    # Check file size
    file.seek(0, os.SEEK_END)
    file_size = file.tell()
    file.seek(0)

    if file_size > MAX_FILE_SIZE:
        return jsonify({'error': f'File too large. Maximum size is {MAX_FILE_SIZE // (1024 * 1024)} MB.'}), 400

    if file_size == 0:
        return jsonify({'error': 'Uploaded file is empty.'}), 400

    try:
        file_bytes = file.read()
        pptx_file = BytesIO(file_bytes)
        presentation = Presentation(pptx_file)

        slides_data = []
        total_word_count = 0

        for i, slide in enumerate(presentation.slides):
            slide_info = extract_slide_text(slide, i)
            slides_data.append(slide_info)
            total_word_count += slide_info['wordCount']

        return jsonify({
            'slideCount': len(slides_data),
            'totalWordCount': total_word_count,
            'slides': slides_data,
            'fileName': file.filename,
        })

    except Exception as e:
        traceback.print_exc()
        return jsonify({'error': f'Failed to parse the presentation: {str(e)}'}), 500


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5001))
    debug = os.environ.get('FLASK_DEBUG', '0') == '1'
    app.run(host='0.0.0.0', port=port, debug=debug)