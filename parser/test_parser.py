from io import BytesIO
from pptx import Presentation
from app import app
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
buf = BytesIO()
prs.save(buf)
buf.seek(0)
with app.test_client() as c:
    data = {'file': (buf, 'test.pptx')}
    resp = c.post('/parse', data=data, content_type='multipart/form-data')
    print('--- Test Results ---')
    print('Status:', resp.status_code)
    print('JSON:', resp.get_json())
